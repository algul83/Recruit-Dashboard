"""PDF / PPTX / HTML / ZIP / 이미지 처리.

- 텍스트 추출: PDF, PPTX, HTML
- Vision 자료 수집: 이미지(.jpg/.png/.webp/.gif), 이미지 기반 PDF, zip 내부 파일
"""
from __future__ import annotations

import io
import re
import zipfile

URL_PATTERN = re.compile(r'https?://[^\s"\'<>()\[\]{}]+')
URL_FETCH_TIMEOUT = 6
URL_FETCH_MAX_BYTES = 10 * 1024 * 1024  # 10MB
MAX_URLS_PER_DOC = 5  # 한 문서당 최대 추출 URL 수
URL_USER_AGENT = (
    'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
    'AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0 Safari/537.36'
)
# fetch 제외할 noise 도메인 (이력서·자소서에 흔히 나오는 의미 없는 링크)
URL_NOISE_PATTERNS = (
    'saramin.co.kr', 'jobkorea.co.kr', 'wanted.co.kr',
    'naver.com/help', 'google.com/help', 'google.com/maps',
    'mailto:', 'tel:', 'fonts.googleapis.com',
)

IMAGE_EXTS = ('.jpg', '.jpeg', '.png', '.gif', '.webp')
IMAGE_MIME = {
    '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
    '.png': 'image/png', '.gif': 'image/gif', '.webp': 'image/webp',
}
MAX_IMAGE_SIZE = 5 * 1024 * 1024  # 5 MB (Anthropic 제한)
MAX_PDF_SIZE = 32 * 1024 * 1024   # 32 MB (Anthropic 제한)
# Anthropic 이미지 차원 한도 — 단일 8000×8000, 다수 1568×1568 권장
# 한 변이 너무 길면 (배너 등) 자동 축소
MAX_IMAGE_DIM = 7500   # 한 변 최대 px
TARGET_LONG_EDGE = 2000  # 축소 시 긴 변 목표 px


def is_image_filename(name: str) -> bool:
    return name.lower().endswith(IMAGE_EXTS)


def normalize_image(data: bytes, media_type: str) -> tuple[bytes, str] | None:
    """이미지가 Anthropic 차원 한도 초과면 PIL로 축소. 정상이면 원본 그대로.

    Returns: (data, media_type) 또는 None (불러올 수 없는 이미지)
    """
    try:
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(data))
        w, h = img.size
        max_dim = max(w, h)
        if max_dim <= MAX_IMAGE_DIM:
            return data, media_type
        # 비율 유지하며 긴 변을 TARGET_LONG_EDGE로 축소
        scale = TARGET_LONG_EDGE / max_dim
        new_size = (max(1, int(w * scale)), max(1, int(h * scale)))
        img = img.convert("RGB") if img.mode in ("P", "RGBA", "LA") else img
        img = img.resize(new_size, Image.LANCZOS)
        buf = io.BytesIO()
        # png는 무손실 + 큼 → jpeg로 통일해 사이즈 축소
        img.save(buf, format="JPEG", quality=85, optimize=True)
        return buf.getvalue(), "image/jpeg"
    except Exception:
        return None


def image_media_type(name: str) -> str:
    for ext, mime in IMAGE_MIME.items():
        if name.lower().endswith(ext):
            return mime
    return 'image/png'


def extract_urls(text: str) -> list[str]:
    """텍스트에서 URL 추출 → noise 제거 + 중복 제거."""
    if not text:
        return []
    raw = URL_PATTERN.findall(text)
    seen: set = set()
    out: list[str] = []
    for u in raw:
        # 끝의 . , ) ] 같은 구두점 제거
        u = u.rstrip('.,);]>')
        ul = u.lower()
        if any(p in ul for p in URL_NOISE_PATTERNS):
            continue
        if u in seen:
            continue
        seen.add(u)
        out.append(u)
        if len(out) >= MAX_URLS_PER_DOC:
            break
    return out


def fetch_url_content(url: str) -> dict | None:
    """URL → {type, data, name, media_type?} 또는 None.

    - text/html → type='text', 본문 텍스트 (≤30000자)
    - PDF → type='pdf', bytes (≤MAX_PDF_SIZE)
    - 이미지 → type='image', bytes (≤MAX_IMAGE_SIZE)
    - 인증 필요·timeout·error → None
    """
    import requests
    from urllib.parse import urlparse
    try:
        r = requests.get(
            url, timeout=URL_FETCH_TIMEOUT, stream=True,
            headers={'User-Agent': URL_USER_AGENT}, allow_redirects=True,
        )
        if r.status_code != 200:
            return None
        ct = (r.headers.get('content-type', '') or '').lower()
        name = urlparse(url).path or url
        # streaming download with size limit
        buf = bytearray()
        for chunk in r.iter_content(chunk_size=64 * 1024):
            if chunk:
                buf.extend(chunk)
                if len(buf) > URL_FETCH_MAX_BYTES:
                    return None
        content = bytes(buf)

        if 'pdf' in ct or name.lower().endswith('.pdf'):
            if len(content) > MAX_PDF_SIZE:
                return None
            return {'type': 'pdf', 'data': content, 'name': f"URL→{name[:60]}"}
        if 'image/' in ct:
            if len(content) > MAX_IMAGE_SIZE:
                return None
            media = ct.split(';')[0].strip()
            return {'type': 'image', 'data': content, 'name': f"URL→{name[:60]}",
                    'media_type': media}
        if 'html' in ct or 'text' in ct or not ct:
            try:
                from bs4 import BeautifulSoup
                html = content.decode('utf-8', errors='replace')
                soup = BeautifulSoup(html, 'html.parser')
                for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
                    tag.decompose()
                text = soup.get_text(separator='\n')
                text = re.sub(r'\n\s*\n', '\n', text).strip()
                if len(text) < 100:
                    return None  # 로그인 페이지 등 빈 페이지
                return {'type': 'text', 'data': text[:30000], 'name': f"URL→{url[:80]}"}
            except Exception:
                return None
        return None
    except Exception:
        return None


def extract_zip(data: bytes) -> list[dict]:
    """zip 풀어서 안의 PDF/이미지/PPTX/HTML 반환 [{name, data, ext}]."""
    out = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for info in zf.infolist():
                if info.is_dir():
                    continue
                # __MACOSX 등 시스템 파일 제외
                if info.filename.startswith('__MACOSX') or info.filename.startswith('.'):
                    continue
                name_lower = info.filename.lower()
                if name_lower.endswith(('.pdf', '.pptx', '.html', '.htm') + IMAGE_EXTS):
                    if info.file_size > MAX_PDF_SIZE:
                        continue
                    content = zf.read(info.filename)
                    ext = name_lower.split('.')[-1]
                    out.append({'name': info.filename, 'data': content, 'ext': ext})
    except Exception:
        pass
    return out


def extract_pdf(data: bytes) -> str:
    """PDF 바이너리 → 텍스트."""
    from pypdf import PdfReader
    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n".join((p.extract_text() or "") for p in reader.pages).strip()
    except Exception as e:
        return f"[PDF 추출 실패: {e}]"


def extract_pptx(data: bytes) -> str:
    """PPTX 바이너리 → 슬라이드별 텍스트."""
    from pptx import Presentation
    try:
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    parts.append(shape.text)
            if parts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))
        return "\n\n".join(slides_text)
    except Exception as e:
        return f"[PPTX 추출 실패: {e}]"


def extract_html(data: bytes) -> str:
    """HTML 바이너리 → 텍스트 (사람인 URL.html 등)."""
    from bs4 import BeautifulSoup
    try:
        html = data.decode('utf-8', errors='replace')
        soup = BeautifulSoup(html, 'html.parser')
        for s in soup(['script', 'style']):
            s.decompose()
        text = soup.get_text(separator='\n')
        # URL 추출 (사람인 첨부 URL이 보통 안에 있음)
        urls = list(set(re.findall(r'https?://[^\s"\'<>]+', html)))
        result = re.sub(r'\n\s*\n', '\n', text).strip()
        if urls:
            result += "\n\n[추출된 URL]\n" + "\n".join(urls)
        return result
    except Exception as e:
        return f"[HTML 추출 실패: {e}]"


def extract_any(filename: str, data: bytes) -> str:
    """파일명 확장자로 적절한 추출기 선택."""
    fname = filename.lower()
    if fname.endswith('.pdf'):
        return extract_pdf(data)
    if fname.endswith('.pptx'):
        return extract_pptx(data)
    if fname.endswith('.html') or fname.endswith('.htm'):
        return extract_html(data)
    return f"[지원 안 되는 형식: {filename}]"
